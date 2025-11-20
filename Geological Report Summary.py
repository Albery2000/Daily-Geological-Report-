# geological_report_upload_app.py
import streamlit as st
import pandas as pd
import numpy as np
from datetime import datetime
import io
import re
from PIL import Image, ImageOps, ImageDraw
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    HAVE_PPTX = True
except Exception:
    Presentation = None
    Inches = None
    Pt = None
    HAVE_PPTX = False

def display_upload_instructions():
    """Show simple upload instructions when no file is provided."""
    st.header("📁 Upload a Daily Geological Report")
    st.markdown(
        """
        Please upload your F-23 style Excel file using the upload control.
        Once uploaded you can view the merged report, formation tops, lithology,
        and gas readings. A PowerPoint summary can be generated when the
        `python-pptx` package is available.
        """
    )
    st.info("Use the left sidebar to upload an `.xlsx` Daily Geological Report file.")
    st.markdown("- If you don't have a file, the app will show demo data.")
    st.markdown("- Images are uploaded inside the 'Merged Report' view.")
    st.markdown("---")

def main():
    # Configure the page
    st.set_page_config(
        page_title="Geological Report Analyzer",
        page_icon="⛰️",
        layout="wide"
    )

    # App title and description
    st.title("🗻 Geological Report Analyzer")
    st.markdown("Upload your Daily Geological Report Excel file to generate automated summaries and analysis.")
    st.markdown("---")

    # File upload section
    uploaded_file = st.sidebar.file_uploader(
        "📤 Upload Excel Report", 
        type=['xlsx'],
        help="Upload your F-23 Daily Geological Report Excel file"
    )

    # (Sidebar image upload removed) — images are uploaded in the 'Merged Report' view

    if uploaded_file is not None:
        try:
            # Load and process the uploaded file
            formation_tops_df, gas_readings_df, lithology_df, well_info, daily_df = process_uploaded_file(uploaded_file)
            
            # Display success message
            st.sidebar.success(f"✅ File loaded successfully: {uploaded_file.name}")
            
            # Navigation
            st.sidebar.title("Navigation")
            section = st.sidebar.radio(
                "Select Section:",
                ["Merged Report", "Report Summary", "Formation Tops", "Lithology Description", "Gas Readings", "Detailed Gas Data", "Raw Data Preview"]
            )

            # Display sections based on selection
            if section == "Report Summary":
                display_report_summary(well_info, formation_tops_df, gas_readings_df)
            elif section == "Formation Tops":
                display_formation_tops(formation_tops_df)
            elif section == "Lithology Description":
                display_lithology_description(lithology_df)
            elif section == "Gas Readings":
                display_gas_readings(gas_readings_df)
            elif section == "Detailed Gas Data":
                display_detailed_gas_data(uploaded_file)
            elif section == "Raw Data Preview":
                display_raw_data_preview(uploaded_file)
            elif section == "Merged Report":
                # Pass image_files explicitly (no sidebar uploader — images uploaded inside Merged Report)
                display_merged_report(formation_tops_df, gas_readings_df, lithology_df, well_info, image_files=None, daily_df=daily_df)

        except Exception as e:
            st.error(f"Error processing file: {str(e)}")
            st.info("Please make sure you're uploading the correct F-23 Daily Geological Report format.")
    else:
        display_upload_instructions()

    # Footer
    display_footer()

def process_uploaded_file(uploaded_file):
    """Process the uploaded Excel file and extract relevant data"""
    
    # Read all sheets
    try:
        # Daily Geological Report sheet
        daily_df = pd.read_excel(uploaded_file, sheet_name='Daily Geological Report', header=None)
        
        # Lithological Description sheet
        lithology_df = pd.read_excel(uploaded_file, sheet_name='Lithological Description', header=None)
        
        # Gas Readings sheet
        gas_df = pd.read_excel(uploaded_file, sheet_name='Lithology %, ROP & Gas Reading', header=None)
        
    except Exception as e:
        st.error(f"Error reading Excel file: {str(e)}")
        # Fallback to demo data
        return get_demo_data()
    
    # Extract well information
    well_info = extract_well_info(daily_df)
    
    # Extract formation tops
    formation_tops_df = extract_formation_tops(daily_df)
    
    # Extract lithology data
    lithology_summary_df = extract_lithology_data(lithology_df)
    
    # Extract gas readings
    gas_readings_df = extract_gas_readings(daily_df)

    return formation_tops_df, gas_readings_df, lithology_summary_df, well_info, daily_df

def extract_well_info(daily_df):
    """Extract well information from the daily report sheet"""
    well_info = {}
    
    try:
        # Find well name - looking for "Well:-" pattern
        for idx, row in daily_df.iterrows():
            row_str = ' '.join([str(cell) for cell in row.values if pd.notna(cell)])
            if 'Well:-' in row_str:
                # Look for well name in nearby cells
                for i, cell in enumerate(row.values):
                    if pd.notna(cell) and 'Ferdaus' in str(cell):
                        well_info["Well Name"] = str(cell)
                        break
                if "Well Name" not in well_info:
                    # Try next row or previous row
                    if idx + 1 < len(daily_df):
                        next_row = daily_df.iloc[idx + 1]
                        for cell in next_row.values:
                            if pd.notna(cell) and 'Ferdaus' in str(cell):
                                well_info["Well Name"] = str(cell)
                                break
        
        # Extract depths at different times
        depth_data = extract_time_depths(daily_df)
        well_info["24:00 Hrs Depth"] = f"{depth_data.get('24:00', 'N/A')} ft"
        well_info["00:00 Hrs Depth"] = f"{depth_data.get('00:00', 'N/A')} ft" 
        well_info["06:00 Hrs Depth"] = f"{depth_data.get('06:00', 'N/A')} ft"
        
        # Calculate progress
        if depth_data.get('24:00') and depth_data.get('00:00'):
            progress_24h = depth_data['00:00'] - depth_data['24:00']
            well_info["Progress (Last 24H)"] = f"{progress_24h} ft"
        
        if depth_data.get('00:00') and depth_data.get('06:00'):
            progress_6h = depth_data['06:00'] - depth_data['00:00']
            well_info["Progress (Last 6H)"] = f"{progress_6h} ft"
        
        # Extract other information
        for idx, row in daily_df.iterrows():
            row_str = ' '.join([str(cell) for cell in row.values if pd.notna(cell)])
            
            if 'Concession:-' in row_str:
                for i, cell in enumerate(row.values):
                    if pd.notna(cell) and 'North Bahariya' in str(cell):
                        well_info["Concession"] = str(cell)
                        break
            
            if 'Date:-' in row_str:
                for i, cell in enumerate(row.values):
                    if pd.notna(cell) and '2025' in str(cell):
                        well_info["Date"] = str(cell)
                        break
            
            if 'Report No.:-' in row_str:
                for i, cell in enumerate(row.values):
                    if pd.notna(cell) and str(cell).isdigit():
                        well_info["Report No."] = str(cell)
                        break
            
            if 'RKB:-' in row_str:
                for i, cell in enumerate(row.values):
                    if pd.notna(cell) and str(cell).isdigit():
                        well_info["RKB"] = f"{cell} ft"
                        break
            
            if 'Spud Date:-' in row_str:
                for i, cell in enumerate(row.values):
                    if pd.notna(cell) and '2025' in str(cell):
                        well_info["Spud Date"] = str(cell)
                        break
            
            if 'Wellsite Geologist:' in row_str:
                # Look for name in nearby cells
                for i in range(idx, min(idx + 3, len(daily_df))):
                    row_check = daily_df.iloc[i]
                    for cell in row_check.values:
                        if pd.notna(cell) and 'Soliman' in str(cell):
                            well_info["Wellsite Geologist"] = str(cell)
                            break
                    if "Wellsite Geologist" in well_info:
                        break
        
        # Set defaults for any missing values
        if "Well Name" not in well_info:
            well_info["Well Name"] = "Ferdaus-23"
        if "Concession" not in well_info:
            well_info["Concession"] = "North Bahariya"
        if "Wellsite Geologist" not in well_info:
            well_info["Wellsite Geologist"] = "Soliman Farag"
            
    except Exception as e:
        st.warning(f"Could not extract all well information: {e}")
        # Return default values
        well_info = get_default_well_info()
    
    return well_info

def extract_time_depths(daily_df):
    """Extract depths at different times from the daily report"""
    depths = {'24:00': None, '00:00': None, '06:00': None}
    
    try:
        for idx, row in daily_df.iterrows():
            for col_idx, cell in enumerate(row.values):
                if pd.isna(cell):
                    continue
                    
                cell_str = str(cell).strip()
                
                # Look for time labels
                if '24:00 Hrs' in cell_str:
                    # Look for depth value in nearby cells
                    for offset in [1, 2, -1, -2]:
                        check_col = col_idx + offset
                        if 0 <= check_col < len(row.values):
                            depth_cell = row.values[check_col]
                            if pd.notna(depth_cell) and is_numeric_value(depth_cell):
                                depths['24:00'] = float(depth_cell)
                                break
                
                elif '00:00 Hrs' in cell_str:
                    for offset in [1, 2, -1, -2]:
                        check_col = col_idx + offset
                        if 0 <= check_col < len(row.values):
                            depth_cell = row.values[check_col]
                            if pd.notna(depth_cell) and is_numeric_value(depth_cell):
                                depths['00:00'] = float(depth_cell)
                                break
                
                elif '06:00 Hrs' in cell_str:
                    for offset in [1, 2, -1, -2]:
                        check_col = col_idx + offset
                        if 0 <= check_col < len(row.values):
                            depth_cell = row.values[check_col]
                            if pd.notna(depth_cell) and is_numeric_value(depth_cell):
                                depths['06:00'] = float(depth_cell)
                                break
        
        # If not found in the expected format, try alternative search
        if not any(depths.values()):
            for idx, row in daily_df.iterrows():
                row_str = ' '.join([str(cell) for cell in row.values if pd.notna(cell)])
                if 'Progress 0-24 Hrs' in row_str or 'Progress Last 6 Hrs' in row_str:
                    # Look in surrounding rows for depth values
                    for row_offset in [-3, -2, -1, 1, 2, 3]:
                        check_idx = idx + row_offset
                        if 0 <= check_idx < len(daily_df):
                            check_row = daily_df.iloc[check_idx]
                            for cell in check_row.values:
                                if pd.notna(cell) and is_numeric_value(cell) and float(cell) > 1000:
                                    if depths['24:00'] is None:
                                        depths['24:00'] = float(cell)
                                    elif depths['00:00'] is None:
                                        depths['00:00'] = float(cell)
                                    elif depths['06:00'] is None:
                                        depths['06:00'] = float(cell)
                                    break
                            
    except Exception as e:
        st.warning(f"Could not extract time depths: {e}")
    
    return depths

def is_numeric_value(value):
    """Check if a value can be converted to float and is a reasonable depth value"""
    try:
        num = float(value)
        return 0 < num < 20000  # Reasonable depth range for oil wells
    except (ValueError, TypeError):
        return False

def extract_formation_tops(daily_df):
    """Extract formation tops data from the actual Excel structure"""
    formation_data = []
    
    try:
        # Find the formation tops table - look for header pattern
        start_row = None
        for idx, row in daily_df.iterrows():
            row_str = ' '.join([str(cell) for cell in row.values if pd.notna(cell)])
            if 'Formation Name' in row_str and 'Member Name' in row_str:
                start_row = idx + 2  # Data starts 2 rows after header
                break
        
        if start_row is not None:
            for idx in range(start_row, min(start_row + 20, len(daily_df))):
                row = daily_df.iloc[idx]
                
                # Extract formation data - based on the Excel structure
                formation_name = safe_get_cell(row, 2)  # Column C
                member_name = safe_get_cell(row, 3)    # Column D
                prognosed_md = safe_get_cell(row, 6)   # Column G
                prognosed_tvdss = safe_get_cell(row, 7) # Column H
                actual_md = safe_get_cell(row, 9)      # Column J
                actual_tvdss = safe_get_cell(row, 10)  # Column K
                
                # Skip empty rows and header-like rows
                if (formation_name and formation_name not in ['Formation Name', 'DABAA', ''] and 
                    not str(formation_name).startswith('Unnamed')):
                    
                    # Build full formation name with member if available
                    full_formation = str(formation_name)
                    if member_name and pd.notna(member_name) and member_name != '':
                        full_formation += f" {member_name}"
                    
                    formation_data.append({
                        'Formation': full_formation,
                        'Prognosed_MD': prognosed_md,
                        'Prognosed_TVDSS': prognosed_tvdss,
                        'Actual_MD': actual_md,
                        'Actual_TVDSS': actual_tvdss
                    })
        
        # If no data found with above method, use fallback
        if not formation_data:
            formation_data = get_fallback_formation_data()
            
    except Exception as e:
        st.warning(f"Could not extract formation tops: {e}")
        formation_data = get_fallback_formation_data()
    
    return pd.DataFrame(formation_data)

def safe_get_cell(row, col_index):
    """Safely get cell value, handling index errors"""
    try:
        if col_index < len(row.values):
            value = row.values[col_index]
            return value if pd.notna(value) else ''
        return ''
    except IndexError:
        return ''

def extract_lithology_data(lithology_df):
    """Extract lithology description data"""
    lithology_data = []
    
    try:
        current_formation = None
        depth_from = None
        depth_to = None
        lithology_desc = []
        
        for idx, row in lithology_df.iterrows():
            row_str = ' '.join([str(cell) for cell in row.values if pd.notna(cell)])
            
            # Look for formation markers
            if any(fm in row_str for fm in ['Moghra Fm.', 'Dabaa Fm.', 'Apollonia Fm', 'Khoman Fm', 
                                           'A/R "', 'Upper Bahariya Fm.']):
                # Save previous formation data if exists
                if current_formation and depth_from and depth_to:
                    lithology_data.append({
                        'Formation': current_formation,
                        'Depth_From': depth_from,
                        'Depth_To': depth_to,
                        'Lithology': ' '.join(lithology_desc)[:200] + '...' if len(' '.join(lithology_desc)) > 200 else ' '.join(lithology_desc)
                    })
                
                # Reset for new formation
                current_formation = extract_formation_name(row_str)
                lithology_desc = []
                
            # Look for depth information
            elif 'Depth:' in row_str and 'From' in row_str:
                depth_from = extract_depth_value(row, lithology_df, idx)
            elif 'To' in row_str and depth_from is not None:
                depth_to = extract_depth_value(row, lithology_df, idx)
                
            # Look for lithology description
            elif 'Lithology:' in row_str or ('*' in row_str and current_formation):
                # Get the lithology description from this row and next rows
                desc = extract_lithology_description(row_str)
                if desc:
                    lithology_desc.append(desc)
        
        # Add the last formation
        if current_formation and depth_from and depth_to:
            lithology_data.append({
                'Formation': current_formation,
                'Depth_From': depth_from,
                'Depth_To': depth_to,
                'Lithology': ' '.join(lithology_desc)[:200] + '...' if len(' '.join(lithology_desc)) > 200 else ' '.join(lithology_desc)
            })
        
    except Exception as e:
        st.warning(f"Could not extract lithology data: {e}")
        lithology_data = get_fallback_lithology_data()
    
    return pd.DataFrame(lithology_data)

def extract_formation_name(row_str):
    """Extract formation name from row string"""
    formations = ['Moghra Fm.', 'Dabaa Fm.', 'Apollonia Fm', 'Khoman Fm', 
                 'A/R "A"', 'A/R "B"', 'A/R "C"', 'A/R "D"', 'A/R "E"', 'A/R "F"',
                 'Upper A/R "G"', 'Middle A/R "G"', 'Lower A/R "G"', 'Upper Bahariya Fm.']
    
    for fm in formations:
        if fm in row_str:
            return fm
    return "Unknown Formation"

def extract_depth_value(row, df, row_idx):
    """Extract depth value from row"""
    for cell in row.values:
        if pd.notna(cell) and is_numeric_value(cell):
            return float(cell)
    
    # Check next row if current row doesn't have numeric value
    if row_idx + 1 < len(df):
        next_row = df.iloc[row_idx + 1]
        for cell in next_row.values:
            if pd.notna(cell) and is_numeric_value(cell):
                return float(cell)
    
    return None

def extract_lithology_description(row_str):
    """Extract lithology description from row string"""
    # Remove common prefixes and clean up
    clean_str = row_str.replace('Lithology:', '').replace('*', '').strip()
    if clean_str and len(clean_str) > 5:  # Meaningful description
        return clean_str
    return ''

def extract_gas_readings(daily_df):
    """Extract gas readings data"""
    gas_data = []
    
    try:
        # Look for gas reading sections in the daily report
        current_zone = None
        max_gas_data = {}
        background_data = {}
        
        for idx, row in daily_df.iterrows():
            row_str = ' '.join([str(cell) for cell in row.values if pd.notna(cell)])
            
            # Look for zone headers
            if 'Max. Gas Reading at:' in row_str:
                # Save previous zone data if exists
                if current_zone and max_gas_data:
                    gas_entry = {
                        'Zone': current_zone,
                        'Max_Gas_Depth': max_gas_data.get('depth'),
                        'TG_Max': max_gas_data.get('TG'),
                        'C1_Max': max_gas_data.get('C1'),
                        'C2_Max': max_gas_data.get('C2'),
                        'C3_Max': max_gas_data.get('C3'),
                        'C4I_Max': max_gas_data.get('C4I'),
                        'C4N_Max': max_gas_data.get('C4N'),
                        'C5_Max': max_gas_data.get('C5')
                    }
                    gas_data.append(gas_entry)
                
                # Reset for new zone
                current_zone = extract_zone_name(row_str, daily_df, idx)
                max_gas_data = {}
                background_data = {}
            
            # Look for depth information for max gas
            elif current_zone and is_numeric_value(row_str) and float(row_str) > 1000:
                max_gas_data['depth'] = float(row_str)
            
            # Look for gas composition data
            elif 'T.G' in row_str or 'C1' in row_str:
                # This might be a header row, check next row for values
                if idx + 1 < len(daily_df):
                    next_row = daily_df.iloc[idx + 1]
                    gas_values = extract_gas_values(next_row)
                    if gas_values:
                        if 'Max' in row_str or not background_data:
                            max_gas_data.update(gas_values)
                        else:
                            background_data.update(gas_values)
        
        # Add the last zone
        if current_zone and max_gas_data:
            gas_entry = {
                'Zone': current_zone,
                'Max_Gas_Depth': max_gas_data.get('depth'),
                'TG_Max': max_gas_data.get('TG'),
                'C1_Max': max_gas_data.get('C1'),
                'C2_Max': max_gas_data.get('C2'),
                'C3_Max': max_gas_data.get('C3'),
                'C4I_Max': max_gas_data.get('C4I'),
                'C4N_Max': max_gas_data.get('C4N'),
                'C5_Max': max_gas_data.get('C5')
            }
            gas_data.append(gas_entry)
        
    except Exception as e:
        st.warning(f"Could not extract gas readings: {e}")
        gas_data = get_fallback_gas_data()
    
    return pd.DataFrame(gas_data)

def extract_zone_name(row_str, df, row_idx):
    """Extract zone name from gas reading section"""
    # Look for zone name in current and following rows
    zones = ['Khoman', 'F', 'UG', 'MG', 'LG', 'Upper Baharyia']
    
    # Check current row
    for zone in zones:
        if zone in row_str:
            return zone
    
    # Check next few rows
    for offset in [1, 2, 3]:
        if row_idx + offset < len(df):
            next_row = df.iloc[row_idx + offset]
            next_str = ' '.join([str(cell) for cell in next_row.values if pd.notna(cell)])
            for zone in zones:
                if zone in next_str:
                    return zone
    
    return "Unknown Zone"

def extract_gas_values(row):
    """Extract gas composition values from row"""
    gas_values = {}
    components = ['TG', 'C1', 'C2', 'C3', 'C4I', 'C4N', 'C5']
    
    numeric_values = []
    for cell in row.values:
        if pd.notna(cell) and is_numeric_value(cell):
            numeric_values.append(float(cell))
    
    # Assign values to components based on position
    for i, value in enumerate(numeric_values):
        if i < len(components):
            gas_values[components[i]] = value
    
    return gas_values if gas_values else None

# Fallback data functions (keep these as backup)
def get_fallback_formation_data():
    return [
        {'Formation': 'DABAA', 'Prognosed_MD': 1221, 'Prognosed_TVDSS': 513, 'Actual_MD': 1216, 'Actual_TVDSS': 508},
        {'Formation': 'APOLLONIA', 'Prognosed_MD': 1960, 'Prognosed_TVDSS': 1252, 'Actual_MD': 1976, 'Actual_TVDSS': 1268},
        {'Formation': 'KHOMAN', 'Prognosed_MD': 3711, 'Prognosed_TVDSS': 3003, 'Actual_MD': 3725, 'Actual_TVDSS': 3017},
        {'Formation': 'A/R "A"', 'Prognosed_MD': 6236, 'Prognosed_TVDSS': 5528, 'Actual_MD': 6205, 'Actual_TVDSS': 5496},
        {'Formation': 'A/R "B"', 'Prognosed_MD': 7181, 'Prognosed_TVDSS': 6469, 'Actual_MD': 7127, 'Actual_TVDSS': 6417},
        {'Formation': 'A/R "C"', 'Prognosed_MD': 7642, 'Prognosed_TVDSS': 6908, 'Actual_MD': 7591, 'Actual_TVDSS': 6872},
        {'Formation': 'A/R "D"', 'Prognosed_MD': 7890, 'Prognosed_TVDSS': 7139, 'Actual_MD': 7851, 'Actual_TVDSS': 7113},
        {'Formation': 'A/R "E"', 'Prognosed_MD': 7951, 'Prognosed_TVDSS': 7196, 'Actual_MD': 'Faulted out', 'Actual_TVDSS': ''},
        {'Formation': 'A/R "F"', 'Prognosed_MD': 8045, 'Prognosed_TVDSS': 7284, 'Actual_MD': 8173, 'Actual_TVDSS': 7397},
        {'Formation': 'Upper A/R "G"', 'Prognosed_MD': 8159, 'Prognosed_TVDSS': 7390, 'Actual_MD': 8243, 'Actual_TVDSS': 7459},
        {'Formation': 'Middle A/R "G"', 'Prognosed_MD': 8546, 'Prognosed_TVDSS': 7750, 'Actual_MD': 8520, 'Actual_TVDSS': 7701},
        {'Formation': 'Lower A/R "G"', 'Prognosed_MD': 8765, 'Prognosed_TVDSS': 7954, 'Actual_MD': 8756, 'Actual_TVDSS': 7908},
        {'Formation': 'Upper Bahariya', 'Prognosed_MD': 8982, 'Prognosed_TVDSS': 8156, 'Actual_MD': 8985, 'Actual_TVDSS': 8108}
    ]

def get_fallback_lithology_data():
    return [
        {'Formation': 'Moghra Fm.', 'Depth_From': 70, 'Depth_To': 1035, 'Lithology': 'SD with clay streaks'},
        {'Formation': 'Dabaa Fm.', 'Depth_From': 1035, 'Depth_To': 1320, 'Lithology': 'SH with SD & LST streaks'},
        {'Formation': 'Apollonia Fm', 'Depth_From': 2910, 'Depth_To': 3115, 'Lithology': 'No return due to complete loss'},
        {'Formation': 'Khoman Fm', 'Depth_From': 3765, 'Depth_To': 4920, 'Lithology': 'CHLKY LST'},
        {'Formation': 'A/R "B"', 'Depth_From': 7500, 'Depth_To': 7591, 'Lithology': 'LST with SH streaks'},
        {'Formation': 'A/R "C"', 'Depth_From': 7591, 'Depth_To': 7851, 'Lithology': 'LST with SH, SLTST, SST streaks'},
        {'Formation': 'A/R "D"', 'Depth_From': 7851, 'Depth_To': 8173, 'Lithology': 'LST with SH streaks'},
        {'Formation': 'A/R "F"', 'Depth_From': 8210, 'Depth_To': 8243, 'Lithology': 'LST with SH streak'},
        {'Formation': 'Upper A/R "G"', 'Depth_From': 8243, 'Depth_To': 8520, 'Lithology': 'SH with LST streaks'},
        {'Formation': 'Middle A/R "G"', 'Depth_From': 8520, 'Depth_To': 8756, 'Lithology': 'SH with SLTST, SST, LST streaks'},
        {'Formation': 'Lower A/R "G"', 'Depth_From': 8756, 'Depth_To': 8985, 'Lithology': 'SH with SLTST, SST, LST streaks'},
        {'Formation': 'Upper Bahariya Fm.', 'Depth_From': 8991, 'Depth_To': 9460, 'Lithology': 'SLTST with SST, SH and LST Streaks'}
    ]

def get_fallback_gas_data():
    return [
        {'Zone': 'Khoman', 'Max_Gas_Depth': 8213, 'TG_Max': 0, 'C1_Max': 0, 'C2_Max': 0, 'C3_Max': 0, 'C4I_Max': 0, 'C4N_Max': 0, 'C5_Max': 0},
        {'Zone': 'F', 'Max_Gas_Depth': 8213, 'TG_Max': 5529, 'C1_Max': 4119, 'C2_Max': 184, 'C3_Max': 3, 'C4I_Max': 0, 'C4N_Max': 0, 'C5_Max': 0},
        {'Zone': 'UG', 'Max_Gas_Depth': 8390, 'TG_Max': 2373, 'C1_Max': 1815, 'C2_Max': 145, 'C3_Max': 66, 'C4I_Max': 40, 'C4N_Max': 10, 'C5_Max': 0},
        {'Zone': 'MG', 'Max_Gas_Depth': 8529, 'TG_Max': 26255, 'C1_Max': 15955, 'C2_Max': 2974, 'C3_Max': 1956, 'C4I_Max': 451, 'C4N_Max': 656, 'C5_Max': 159},
        {'Zone': 'LG', 'Max_Gas_Depth': 8796, 'TG_Max': 137619, 'C1_Max': 77029, 'C2_Max': 15269, 'C3_Max': 7763, 'C4I_Max': 1900, 'C4N_Max': 1910, 'C5_Max': 1020}
    ]

def get_default_well_info():
    return {
        "Well Name": "Ferdaus-23",
        "Concession": "North Bahariya",
        "Date": "2025-06-04",
        "Report No.": "16",
        "RKB": "708 ft",
        "Spud Date": "2025-05-18",
        "24:00 Hrs Depth": "8996 ft",
        "00:00 Hrs Depth": "9460 ft", 
        "06:00 Hrs Depth": "9460 ft",
        "Progress (Last 24H)": "464 ft",
        "Progress (Last 6H)": "0 ft",
        "Wellsite Geologist": "Soliman Farag"
    }

def get_demo_data():
    formation_tops_df = pd.DataFrame(get_fallback_formation_data())
    gas_readings_df = pd.DataFrame(get_fallback_gas_data())
    lithology_df = pd.DataFrame(get_fallback_lithology_data())
    well_info = get_default_well_info()
    
    # Create sample daily_df
    demo_daily = pd.DataFrame([
        ["24:00 Hrs", 8996],
        ["00:00 Hrs", 9460], 
        ["06:00 Hrs", 9460]
    ])
    
    return formation_tops_df, gas_readings_df, lithology_df, well_info, demo_daily

# Image Editor removed — images are uploaded and handled inside the Merged Report view now.

def display_report_summary(well_info, formation_tops_df, gas_readings_df):
    st.header("📊 Report Summary")
    
    # Display well information in columns
    col1, col2, col3 = st.columns(3)
    
    with col1:
        st.subheader("Well Information")
        for key, value in list(well_info.items())[:4]:
            st.write(f"**{key}:** {value}")
    
    with col2:
        st.subheader("Drilling Progress")
        for key, value in list(well_info.items())[4:9]:
            st.write(f"**{key}:** {value}")
    
    with col3:
        st.subheader("Personnel")
        for key, value in list(well_info.items())[9:]:
            st.write(f"**{key}:** {value}")
    
    st.markdown("---")
    
    # Key highlights
    st.subheader("🔍 Key Highlights")
    
    col1, col2 = st.columns(2)
    
    with col1:
        st.info("**Formation Tops**")
        current_formation = formation_tops_df['Formation'].iloc[-1]
        st.write(f"- Currently drilling in **{current_formation}**")
        
        # Count faulted formations
        faulted_count = len(formation_tops_df[formation_tops_df['Actual_MD'] == 'Faulted out'])
        if faulted_count > 0:
            st.write(f"- {faulted_count} formation(s) faulted out")
        
        # Show depth comparison for current formation
        current_row = formation_tops_df.iloc[-1]
        if current_row['Actual_MD'] != 'Faulted out' and is_numeric_value(current_row['Actual_MD']):
            actual_depth = float(current_row['Actual_MD'])
            prognosed_depth = float(current_row['Prognosed_MD'])
            difference = actual_depth - prognosed_depth
            st.write(f"- Current depth: {actual_depth} ft MD")
            st.write(f"- Prognosed depth: {prognosed_depth} ft MD")
            st.write(f"- Difference: {difference:+.1f} ft")
    
    with col2:
        st.warning("**Gas Readings**")
        if not gas_readings_df.empty:
            max_gas_zone = gas_readings_df.loc[gas_readings_df['TG_Max'].idxmax()]
            st.write(f"- Maximum gas: **{max_gas_zone['TG_Max']:,} ppm TG**")
            st.write(f"- At depth: **{max_gas_zone['Max_Gas_Depth']} ft**")
            st.write(f"- Zone: **{max_gas_zone['Zone']}**")
            
            # Count zones with significant gas
            significant_gas = len(gas_readings_df[gas_readings_df['TG_Max'] > 1000])
            st.write(f"- {significant_gas} zone(s) with significant gas shows")
        
        st.success("**Operations**")
        st.write("- Continuous coring operations")
        st.write("- Detailed lithological analysis")
        st.write("- Regular gas monitoring")

def display_formation_tops(formation_tops_df):
    st.header("🏔️ Formation Tops Correlation")
    
    # Create a display version with differences calculated
    display_df = formation_tops_df.copy()
    
    # Calculate depth differences where possible
    differences = []
    for _, row in display_df.iterrows():
        if (is_numeric_value(row['Actual_MD']) and 
            is_numeric_value(row['Prognosed_MD']) and
            row['Actual_MD'] != 'Faulted out'):
            actual = float(row['Actual_MD'])
            prognosed = float(row['Prognosed_MD'])
            diff = actual - prognosed
            differences.append(f"{diff:+.1f} ft")
        else:
            differences.append("N/A")
    
    display_df['Difference'] = differences
    
    # Display the formation tops table
    st.dataframe(
        display_df,
        use_container_width=True,
        hide_index=True
    )
    
    st.markdown("---")
    
    # Visualization of formation tops
    st.subheader("Formation Tops Depth Profile")
    
    # Create a simple depth chart for formations with numeric values
    chart_data = formation_tops_df.copy()
    chart_data = chart_data[chart_data['Actual_MD'] != 'Faulted out']
    chart_data['Actual_MD'] = pd.to_numeric(chart_data['Actual_MD'], errors='coerce')
    chart_data = chart_data.dropna()
    
    if not chart_data.empty:
        # Create a horizontal bar chart
        chart_data = chart_data.sort_values('Actual_MD', ascending=False)
        st.bar_chart(chart_data.set_index('Formation')['Actual_MD'])
    
    # Additional information
    col1, col2 = st.columns(2)
    
    with col1:
        st.info("**Key Formation Details**")
        for _, row in formation_tops_df.iterrows():
            if row['Actual_MD'] != 'Faulted out':
                st.write(f"**{row['Formation']}:** {row['Actual_MD']} ft MD")
            else:
                st.write(f"**{row['Formation']}:** {row['Actual_MD']}")
    
    with col2:
        st.info("**Geological Remarks**")
        current_formation = formation_tops_df['Formation'].iloc[-1]
        st.write(f"- **Current formation:** {current_formation}")
        
        # Count successful correlations
        good_correlations = len(formation_tops_df[
            (formation_tops_df['Actual_MD'] != 'Faulted out') & 
            (pd.to_numeric(formation_tops_df['Actual_MD'], errors='coerce') - 
             pd.to_numeric(formation_tops_df['Prognosed_MD'], errors='coerce')).abs() < 100
        ])
        st.write(f"- **{good_correlations}/{len(formation_tops_df)}** formations with good correlation")
        
        if 'Faulted out' in formation_tops_df['Actual_MD'].values:
            st.write("- **Structural complexity** observed with faulted sections")

def display_lithology_description(lithology_df):
    st.header("🪨 Lithology Description")
    
    # Display lithology table
    st.dataframe(
        lithology_df,
        use_container_width=True,
        hide_index=True
    )
    
    st.markdown("---")
    
    # Detailed lithology descriptions
    st.subheader("Formation Lithology Details")
    
    for _, row in lithology_df.iterrows():
        with st.expander(f"{row['Formation']} ({row['Depth_From']}'-{row['Depth_To']}')"):
            st.write(f"**Depth Interval:** {row['Depth_From']} - {row['Depth_To']} ft")
            st.write(f"**Lithology:** {row['Lithology']}")
            
            # Add formation-specific notes
            if 'Moghra' in row['Formation']:
                st.info("**Characteristics:** Sandstone dominated with clay interbeds")
            elif 'Dabaa' in row['Formation']:
                st.info("**Characteristics:** Shale dominated with carbonate streaks")
            elif 'Apollonia' in row['Formation']:
                st.warning("**Note:** Lost circulation zones present")
            elif 'Khoman' in row['Formation']:
                st.info("**Characteristics:** Chalky limestone formation")
            elif 'A/R' in row['Formation']:
                st.info("**Characteristics:** Complex mixed carbonate-clastic sequence")
            elif 'Bahariya' in row['Formation']:
                st.success("**Characteristics:** Clastic reservoir section with hydrocarbon shows")
    
    # Summary statistics
    st.subheader("📊 Lithology Summary")
    col1, col2, col3 = st.columns(3)
    
    with col1:
        total_interval = sum(row['Depth_To'] - row['Depth_From'] for _, row in lithology_df.iterrows())
        st.metric("Total Interval Described", f"{total_interval} ft")
    
    with col2:
        formations_count = len(lithology_df)
        st.metric("Formations Described", formations_count)
    
    with col3:
        avg_interval = total_interval / formations_count if formations_count > 0 else 0
        st.metric("Average Interval", f"{avg_interval:.0f} ft")

def display_gas_readings(gas_readings_df):
    st.header("💨 Gas Readings Summary")
    
    # Display gas readings table
    st.dataframe(
        gas_readings_df,
        use_container_width=True,
        hide_index=True
    )
    
    st.markdown("---")
    
    # Gas readings visualization
    col1, col2 = st.columns(2)
    
    with col1:
        st.subheader("Total Gas (TG) by Zone")
        if not gas_readings_df.empty:
            tg_chart_data = gas_readings_df[['Zone', 'TG_Max']].set_index('Zone')
            st.bar_chart(tg_chart_data)
    
    with col2:
        st.subheader("Methane (C1) by Zone")
        if not gas_readings_df.empty:
            c1_chart_data = gas_readings_df[['Zone', 'C1_Max']].set_index('Zone')
            st.bar_chart(c1_chart_data)
    
    # Gas composition analysis
    st.subheader("Gas Composition Analysis")
    
    for _, row in gas_readings_df.iterrows():
        with st.expander(f"Zone: {row['Zone']} - Depth: {row['Max_Gas_Depth']} ft"):
            col1, col2, col3 = st.columns(3)
            
            with col1:
                st.metric("Total Gas", f"{row['TG_Max']:,} ppm")
                st.metric("Methane (C1)", f"{row['C1_Max']:,} ppm")
            
            with col2:
                st.metric("Ethane (C2)", f"{row['C2_Max']:,} ppm")
                st.metric("Propane (C3)", f"{row['C3_Max']:,} ppm")
            
            with col3:
                st.metric("Iso-Butane (C4I)", f"{row['C4I_Max']:,} ppm")
                st.metric("Normal-Butane (C4N)", f"{row['C4N_Max']:,} ppm")
                st.metric("Pentane (C5)", f"{row['C5_Max']:,} ppm")
            
            # Gas ratio analysis
            if row['C1_Max'] > 0:
                wetness_ratio = (row['C2_Max'] + row['C3_Max'] + row['C4I_Max'] + row['C4N_Max'] + row['C5_Max']) / row['C1_Max'] * 100
                st.write(f"**Wetness Ratio:** {wetness_ratio:.2f}%")
                
                if wetness_ratio > 5:
                    st.success("**Interpretation:** Wet gas signature")
                elif wetness_ratio > 1:
                    st.info("**Interpretation:** Mixed gas signature")
                else:
                    st.warning("**Interpretation:** Dry gas signature")

def display_detailed_gas_data(uploaded_file):
    st.header("📈 Detailed Gas Readings Analysis")
    
    try:
        # Try to read the actual gas data from the uploaded file
        gas_df = pd.read_excel(uploaded_file, sheet_name='Lithology %, ROP & Gas Reading')
        
        st.success("✅ Loaded detailed gas data from uploaded file")
        
        # Display preview of the data
        st.subheader("Gas Data Preview")
        st.dataframe(gas_df.head(20), use_container_width=True)
        
        # Basic statistics
        st.subheader("Gas Data Statistics")
        numeric_columns = gas_df.select_dtypes(include=[np.number]).columns
        
        if len(numeric_columns) > 0:
            col1, col2, col3, col4 = st.columns(4)
            with col1:
                st.metric("Data Points", len(gas_df))
            with col2:
                if 'DEPTH' in gas_df.columns:
                    st.metric("Depth Range", f"{gas_df['DEPTH'].min():.0f}-{gas_df['DEPTH'].max():.0f} ft")
            with col3:
                if 'TG' in gas_df.columns:
                    st.metric("Max TG", f"{gas_df['TG'].max():,} ppm")
            with col4:
                if 'TG' in gas_df.columns:
                    st.metric("Avg TG", f"{gas_df['TG'].mean():.0f} ppm")
        
        # Depth-based gas plot if depth data available
        if 'DEPTH' in gas_df.columns and 'TG' in gas_df.columns:
            st.subheader("Gas vs Depth")
            plot_data = gas_df[['DEPTH', 'TG']].dropna()
            if not plot_data.empty:
                st.line_chart(plot_data.set_index('DEPTH'))
        
    except Exception as e:
        st.warning("Could not load detailed gas data from the file.")
        st.info("This section requires the 'Lithology %, ROP & Gas Reading' sheet with detailed gas data.")

def display_raw_data_preview(uploaded_file):
    st.header("📋 Raw Data Preview")
    
    sheet_names = ['Daily Geological Report', 'Lithological Description', 'Lithology %, ROP & Gas Reading']
    
    selected_sheet = st.selectbox("Select Sheet to Preview", sheet_names)
    
    try:
        df = pd.read_excel(uploaded_file, sheet_name=selected_sheet)
        st.subheader(f"Sheet: {selected_sheet}")
        
        # Show first 50 rows to avoid overwhelming the interface
        st.dataframe(df.head(50), use_container_width=True)
        
        if len(df) > 50:
            st.info(f"Showing first 50 of {len(df)} rows. Use download for full data.")
        
        # Show basic info about the sheet
        col1, col2, col3 = st.columns(3)
        with col1:
            st.metric("Rows", len(df))
        with col2:
            st.metric("Columns", len(df.columns))
        with col3:
            non_empty_cells = df.count().sum()
            total_cells = len(df) * len(df.columns)
            st.metric("Data Density", f"{(non_empty_cells/total_cells*100):.1f}%")
            
        # Download option
        csv = df.to_csv(index=False)
        st.download_button(
            label="📥 Download CSV",
            data=csv,
            file_name=f"{selected_sheet.replace(' ', '_')}.csv",
            mime="text/csv"
        )
            
    except Exception as e:
        st.error(f"Could not load sheet: {selected_sheet}. Error: {e}")

def display_merged_report(formation_tops_df, gas_readings_df, lithology_df, well_info, image_files, daily_df):
    """Enhanced merged view with well location image marking"""
    st.header("📋 Merged Report - Complete Overview")

    # Well header
    st.subheader(f"🛢️ Well: {well_info.get('Well Name', 'Unknown')}")
    
    # Main information columns
    col1, col2 = st.columns(2)
    
    with col1:
        st.subheader("📊 Drilling Progress")
        st.metric("24:00 Hrs Depth", well_info.get("24:00 Hrs Depth", "N/A"))
        st.metric("00:00 Hrs Depth", well_info.get("00:00 Hrs Depth", "N/A")) 
        st.metric("06:00 Hrs Depth", well_info.get("06:00 Hrs Depth", "N/A"))
        st.metric("Progress (Last 24H)", well_info.get("Progress (Last 24H)", "N/A"))
        st.metric("Progress (Last 6H)", well_info.get("Progress (Last 6H)", "N/A"))
    
    with col2:
        st.subheader("🏢 Well Information")
        st.write(f"**Concession:** {well_info.get('Concession', 'N/A')}")
        st.write(f"**Date:** {well_info.get('Date', 'N/A')}")
        st.write(f"**Report No.:** {well_info.get('Report No.', 'N/A')}")
        st.write(f"**RKB:** {well_info.get('RKB', 'N/A')}")
        st.write(f"**Spud Date:** {well_info.get('Spud Date', 'N/A')}")
        st.write(f"**Geologist:** {well_info.get('Wellsite Geologist', 'N/A')}")

    st.markdown("---")

    # Current formation status
    st.subheader("🎯 Current Drilling Status")
    current_formation = formation_tops_df['Formation'].iloc[-1]
    current_actual_md = formation_tops_df['Actual_MD'].iloc[-1]
    current_prognosed_md = formation_tops_df['Prognosed_MD'].iloc[-1]
    
    col1, col2, col3 = st.columns(3)
    
    with col1:
        st.metric("Current Formation", current_formation)
    
    with col2:
        st.metric("Actual Depth", f"{current_actual_md} ft")
    
    with col3:
        if current_actual_md != 'Faulted out' and is_numeric_value(current_actual_md) and is_numeric_value(current_prognosed_md):
            difference = float(current_actual_md) - float(current_prognosed_md)
            st.metric("vs Prognosed", f"{difference:+.1f} ft")
        else:
            st.metric("vs Prognosed", "N/A")

    st.markdown("---")

    # Formation tops summary
    st.subheader("🏔️ Formation Tops Summary")
    
    # Create display table with differences
    display_tops = formation_tops_df.copy()
    differences = []
    for _, row in display_tops.iterrows():
        if (is_numeric_value(row['Actual_MD']) and 
            is_numeric_value(row['Prognosed_MD']) and
            row['Actual_MD'] != 'Faulted out'):
            diff = float(row['Actual_MD']) - float(row['Prognosed_MD'])
            differences.append(f"{diff:+.1f} ft")
        else:
            differences.append("N/A")
    
    display_tops['Difference'] = differences
    st.dataframe(display_tops, use_container_width=True)

    st.markdown("---")

    # Well location images section
    st.subheader("🗺️ Well Location & Site Images")
    
    if image_files:
        st.info("💡 Use the Image Editor tab to mark well locations on these images")
        
        # Display available images
        cols = st.columns(min(3, len(image_files)))
        for idx, img_file in enumerate(image_files):
            with cols[idx % 3]:
                try:
                    img = Image.open(img_file)
                    st.image(img, use_column_width=True, caption=img_file.name)
                except Exception as e:
                    st.error(f"Could not display {img_file.name}")
    else:
        st.warning("No well location images uploaded. Use the sidebar to upload map or site photos.")

    st.markdown("---")

    # Gas readings summary
    st.subheader("💨 Gas Readings Summary")
    if not gas_readings_df.empty:
        st.dataframe(gas_readings_df, use_container_width=True)
        
        # Highlight maximum gas reading
        max_gas_row = gas_readings_df.loc[gas_readings_df['TG_Max'].idxmax()]
        st.warning(f"🚨 **Maximum Gas Reading:** {max_gas_row['TG_Max']:,} ppm TG in {max_gas_row['Zone']} zone at {max_gas_row['Max_Gas_Depth']} ft")
    else:
        st.info("No gas reading data available")

    st.markdown("---")

    # Export options
    st.subheader("📤 Export Report")

    # Detect python-pptx at runtime so UI can show the button without a full server restart
    try:
        from pptx import Presentation  # noqa: F401
        have_pptx = True
    except Exception:
        have_pptx = False

    if have_pptx:
        if st.button("🎯 Generate PowerPoint Summary Report"):
            with st.spinner("Creating PowerPoint presentation..."):
                # Get current depths from well_info
                depths = {}
                for time_key in ['24:00 Hrs Depth', '00:00 Hrs Depth', '06:00 Hrs Depth']:
                    depth_str = well_info.get(time_key, '0 ft')
                    depth_value = float(''.join(filter(str.isdigit, depth_str))) if 'ft' in depth_str else 0
                    depths[time_key.replace(' Hrs Depth', '')] = depth_value

                progress = {
                    '0-24': well_info.get('Progress (Last 24H)', ''),
                    '6': well_info.get('Progress (Last 6H)', '')
                }

                pptx_bytes = create_presentation(
                    formation_tops_df, lithology_df, gas_readings_df,
                    well_info, None, None, depths=depths, progress=progress
                )

                if pptx_bytes:
                    st.success("✅ PowerPoint generated successfully!")
                    st.download_button(
                        "📥 Download PPTX Report",
                        data=pptx_bytes.getvalue(),
                        file_name=f"{well_info.get('Well Name','well')}_geological_report.pptx",
                        mime='application/vnd.openxmlformats-officedocument.presentationml.presentation'
                    )
    else:
        st.warning("📝 PowerPoint export requires 'python-pptx' package. Install it to enable export feature.")

def create_presentation(formation_tops_df, lithology_df, gas_readings_df, well_info, img1, img2, depths=None, progress=None):
    """Create PowerPoint presentation - simplified version"""
    try:
        prs = Presentation()

        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = f"Geological Report - {well_info.get('Well Name','Unknown')}"
        subtitle.text = f"Date: {well_info.get('Date','Unknown')} | Report No: {well_info.get('Report No.','')}"

        # Well summary slide
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Well Summary"
        content = slide.shapes.placeholders[1].text_frame
        content.clear()
        
        summary_points = [
            f"Well: {well_info.get('Well Name','')}",
            f"Concession: {well_info.get('Concession','')}",
            f"Current Depth: {well_info.get('06:00 Hrs Depth','')}",
            f"24H Progress: {well_info.get('Progress (Last 24H)','')}",
            f"6H Progress: {well_info.get('Progress (Last 6H)','')}",
            f"Geologist: {well_info.get('Wellsite Geologist','')}"
        ]
        
        for point in summary_points:
            p = content.add_paragraph()
            p.text = point

        # Formation tops slide
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Formation Tops"
        content = slide.shapes.placeholders[1].text_frame
        content.clear()
        
        # Add key formation tops
        for _, row in formation_tops_df.iterrows():
            p = content.add_paragraph()
            p.text = f"{row['Formation']}: Actual {row['Actual_MD']} ft | Prognosed {row['Prognosed_MD']} ft"

        # Gas readings slide
        if not gas_readings_df.empty:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Gas Readings Summary"
            content = slide.shapes.placeholders[1].text_frame
            content.clear()
            
            for _, row in gas_readings_df.iterrows():
                p = content.add_paragraph()
                p.text = f"{row['Zone']}: {row['TG_Max']:,} ppm TG at {row['Max_Gas_Depth']} ft"

        # Save to BytesIO
        out = io.BytesIO()
        prs.save(out)
        out.seek(0)
        return out
        
    except Exception as e:
        st.error(f"Error creating presentation: {e}")
        return None

def display_footer():
    st.markdown("---")
    st.markdown(
        "**Geological Report Analyzer** | North Bahariya Petroleum Company | "
        "Generated on {}".format(datetime.now().strftime("%Y-%m-%d %H:%M"))
    )

if __name__ == "__main__":
    main()

